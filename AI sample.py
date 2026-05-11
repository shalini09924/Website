import streamlit as st
st.set_page_config(page_title="AI Lecture Assistant", page_icon="🧠", layout="wide")

import os
import whisper
import sounddevice as sd
import tempfile
import scipy.io.wavfile as wav
import json
import time
import re
import contextlib
import threading
import queue
import numpy as np
import torch
from datetime import datetime
from sqlite3 import connect as sqlite_connect
from PIL import Image
import base64
import io
import pandas as pd

try:
    import cv2
    import mss
    VIDEO_RECORDING_AVAILABLE = True
except ImportError:
    VIDEO_RECORDING_AVAILABLE = False

try:
    import PyPDF2
except ImportError:
    PyPDF2 = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

try:
    import ollama
    OLLAMA_AVAILABLE = True
except ImportError:
    OLLAMA_AVAILABLE = False

try:
    from groq import Groq as GroqClient
    GROQ_AVAILABLE = True
except ImportError:
    GROQ_AVAILABLE = False

try:
    from huggingface_hub import InferenceClient as HFInferenceClient
    HF_AVAILABLE = True
except ImportError:
    HF_AVAILABLE = False

try:
    import pytesseract
    import shutil
    tesseract_path = shutil.which("tesseract")
    if tesseract_path:
        pytesseract.pytesseract.tesseract_cmd = tesseract_path
    elif os.path.exists(r"C:\Program Files\Tesseract-OCR\tesseract.exe"):
        pytesseract.pytesseract.tesseract_cmd = r"C:\Program Files\Tesseract-OCR\tesseract.exe"
    OCR_AVAILABLE = True
except ImportError:
    OCR_AVAILABLE = False

# =========================================================
# CONSTANTS AND CONFIG
# =========================================================
BASE_DIR = "app_data"
REP_DIR = os.path.join(BASE_DIR, "reports")
ATTENDANCE_FILE = os.path.join(BASE_DIR, "attendance.csv")
os.makedirs(REP_DIR, exist_ok=True)

# ---------- AI Model Settings ----------
OLLAMA_MODEL = "llama3"
GROQ_MODEL = "llama-3.3-70b-versatile"
GROQ_BACKUPS = ["qwen-2.5-32b", "deepseek-r1-distill-llama-70b"]
HF_MODEL = "mistralai/Mistral-7B-Instruct-v0.3"
HF_BACKUPS = ["Qwen/Qwen2.5-7B-Instruct", "meta-llama/Llama-3.3-70B-Instruct"]

AUDIO_SAMPLE_RATE = 16000
SCREEN_CAPTURE_WIDTH = 800
SCREEN_CAPTURE_HEIGHT = 600
MAX_FILE_SIZE = 50 * 1024 * 1024  # 50MB

# =========================================================
# SESSION STATE INITIALIZATION
# =========================================================
def init_session_state():
    defaults = {
        "running": False,
        "transcript": "",
        "summary": "",
        "last_summary_len": 0,
        "error_msg": "",
        "uploaded_materials": set(),
        "material_text": "",
        "student_questions": [],
        "last_vision_time": 0,
        "recording_stop_event": threading.Event(),
        "ai_provider": "Ollama",
        "groq_client": None,
        "hf_client": None,
        "ollama_model": OLLAMA_MODEL,
        "attendance_data": [],
        "audio_queue": queue.Queue(),
        "show_screen": True,
        "record_screen": False,
        "live_audio": True,
        "silence_filter": 5,
        "summary_trigger": 400,
        "recording_duration": 5,
    }
    for key, default_value in defaults.items():
        if key not in st.session_state:
            st.session_state[key] = default_value

init_session_state()

st.title("🧠 AI Lecture Assistant — Enhanced Edition")

if st.session_state.error_msg:
    st.warning(st.session_state.error_msg)
    st.session_state.error_msg = ""

# =========================================================
# DATABASE HELPERS (Notes)
# =========================================================
@contextlib.contextmanager
def get_db_connection():
    db_path = os.path.join(BASE_DIR, "notes.db")
    conn = sqlite_connect(db_path)
    try:
        yield conn
    finally:
        conn.close()

def init_notes_db():
    try:
        with get_db_connection() as conn:
            conn.execute(
                "CREATE TABLE IF NOT EXISTS notes ("
                "id INTEGER PRIMARY KEY AUTOINCREMENT, "
                "title TEXT NOT NULL, "
                "content TEXT NOT NULL, "
                "created_at TEXT NOT NULL)"
            )
            conn.commit()
    except Exception as e:
        st.session_state.error_msg = f"Database init error: {e}"

def add_note(title, content):
    try:
        with get_db_connection() as conn:
            conn.execute(
                "INSERT INTO notes (title, content, created_at) VALUES (?, ?, ?)",
                (title, content, datetime.now().isoformat())
            )
            conn.commit()
        return True
    except Exception as e:
        st.session_state.error_msg = f"Note save error: {e}"
        return False

def search_notes(query):
    try:
        with get_db_connection() as conn:
            cursor = conn.execute(
                "SELECT title, content, created_at FROM notes WHERE title LIKE ? OR content LIKE ?",
                (f"%{query}%", f"%{query}%")
            )
            return cursor.fetchall()
    except Exception as e:
        st.session_state.error_msg = f"Search error: {e}"
        return []

init_notes_db()

# ---------- Attendance helpers ----------
def init_attendance_file():
    if not os.path.exists(ATTENDANCE_FILE):
        df = pd.DataFrame(columns=["Date", "Time", "Name", "Status", "AI_Support"])
        df.to_csv(ATTENDANCE_FILE, index=False)

def mark_attendance(name, status="On-Time"):
    try:
        if not name or not name.strip():
            st.session_state.error_msg = "Student name cannot be empty"
            return False
        
        now = datetime.now()
        support = "AI Summary Required" if status == "Late" else "Optional"
        data = {
            "Date": [now.strftime("%Y-%m-%d")],
            "Time": [now.strftime("%H:%M:%S")],
            "Name": [name.strip()],
            "Status": [status],
            "AI_Support": [support]
        }
        df = pd.DataFrame(data)
        df.to_csv(ATTENDANCE_FILE, mode='a', index=False, header=False)
        return True
    except Exception as e:
        st.session_state.error_msg = f"Attendance error: {e}"
        return False

def get_latecomers():
    init_attendance_file()
    try:
        df = pd.read_csv(ATTENDANCE_FILE)
        today = datetime.now().strftime("%Y-%m-%d")
        late_mask = (df['Date'] == today) & (df['Status'] == 'Late')
        return df[late_mask]['Name'].unique().tolist()
    except Exception:
        return []

init_attendance_file()

# =========================================================
# UI LAYOUT
# =========================================================
col_screen, col_transcript, col_summary = st.columns([1.2, 1, 1])

with col_screen:
    st.subheader("📺 Presentation View")
    screen_placeholder = st.empty()

with col_transcript:
    st.subheader("📄 Live Transcript")
    transcript_placeholder = st.empty()

with col_summary:
    st.subheader("🧠 AI Summary")
    summary_placeholder = st.empty()

# =========================================================
# SIDEBAR CONTROLS
# =========================================================
with st.sidebar:
    st.subheader("🤖 AI Provider Setup")
    
    ai_provider_choice = st.radio(
        "Select AI Provider",
        ["Ollama (Local, Free)", "Groq (Cloud, Fast)", "Hugging Face (Cloud)"],
        index=0,
        help="Ollama runs locally (no API key). Groq and Hugging Face need free API keys."
    )
    st.session_state.ai_provider = ai_provider_choice.split(" ")[0]
    
    if "Ollama" in ai_provider_choice:
        if not OLLAMA_AVAILABLE:
            st.error("❌ Ollama not installed. Install with: pip install ollama")
            st.info("Then run: ollama serve")
            st.info("Then pull a model: ollama pull llama3")
        else:
            ollama_model_input = st.text_input("Ollama Model", st.session_state.ollama_model,
                                         help="Type any model you have pulled (e.g., deepseek-r1:1.5b)")
            st.session_state.ollama_model = ollama_model_input
            st.caption("You have: llama3, phi4-mini, llama3.1, phi3, qwen2.5:3b, deepseek-r1:1.5b")
    
    elif "Groq" in ai_provider_choice:
        if not GROQ_AVAILABLE:
            st.error("❌ Groq not installed. Install with: pip install groq")
        else:
            groq_key = st.text_input(
                "🔑 Groq API Key",
                type="password",
                placeholder="gsk_...",
                help="Get free key from https://console.groq.com"
            )
            if groq_key:
                try:
                    st.session_state.groq_client = GroqClient(api_key=groq_key)
                    st.success("✅ Groq connected")
                except Exception as e:
                    st.error(f"❌ Groq error: {e}")
                    st.session_state.groq_client = None
            st.caption(f"Primary model: {GROQ_MODEL}")

    elif "Hugging Face" in ai_provider_choice:
        if not HF_AVAILABLE:
            st.error("❌ Hugging Face hub not installed. Install with: pip install huggingface_hub")
        else:
            hf_token = st.text_input(
                "🤗 HF Token",
                type="password",
                placeholder="hf_...",
                help="Get your token from https://huggingface.co/settings/tokens"
            )
            if hf_token:
                try:
                    st.session_state.hf_client = HFInferenceClient(token=hf_token)
                    st.success("✅ Hugging Face connected")
                except Exception as e:
                    st.error(f"❌ HF error: {e}")
                    st.session_state.hf_client = None
            st.caption(f"Primary model: {HF_MODEL}")

    st.divider()
    st.header("⚙️ Session Controls")
    subject = st.text_input("📘 Subject", "CS101").strip().lower().replace(" ", "_")
    session_id = st.text_input("🆔 Session ID", datetime.now().strftime("%Y%m%d_%H%M"))
    topic = st.text_input("📚 Topic", "Lecture Demo")
    
    st.divider()
    st.session_state.live_audio = st.checkbox("🎤 Live Voice Capture", value=st.session_state.live_audio)
    st.session_state.show_screen = st.checkbox("📺 Show Presentation Screen", value=st.session_state.show_screen, disabled=not VIDEO_RECORDING_AVAILABLE)
    st.session_state.record_screen = st.checkbox("🎥 Record Screen to Video", value=st.session_state.record_screen, disabled=not VIDEO_RECORDING_AVAILABLE)
    
    st.session_state.silence_filter = st.slider("🔇 Minimum chars to transcribe", 0, 50, st.session_state.silence_filter,
                               help="Filters out very short/silent utterances")
    st.session_state.summary_trigger = st.slider("Summary after (chars)", 100, 1000, st.session_state.summary_trigger)
    st.session_state.recording_duration = st.slider("🎙️ Recording duration (sec)", 2, 10, st.session_state.recording_duration)
    
    # ----- Attendance -----
    st.divider()
    st.subheader("📋 Attendance")
    enable_attendance = st.checkbox("Track Attendance", value=False)
    if enable_attendance:
        student_name = st.text_input("Student Name (for manual entry)", key="student_name_attendance")
        col_on_time, col_late = st.columns(2)
        with col_on_time:
            if st.button("✅ On-Time", key="on_time_btn"):
                if student_name:
                    if mark_attendance(student_name, "On-Time"):
                        st.success(f"✅ {student_name} marked On-Time")
                else:
                    st.error("Please enter a student name")
        with col_late:
            if st.button("⏰ Late", key="late_btn"):
                if student_name:
                    if mark_attendance(student_name, "Late"):
                        st.success(f"⏰ {student_name} marked Late")
                else:
                    st.error("Please enter a student name")
        
        latecomers = get_latecomers()
        if latecomers:
            st.warning(f"🤖 Late students: {', '.join(latecomers)}")

    # ----- Quick Screen Summary -----
    st.divider()
    st.subheader("🖥️ Quick Screen Summary")
    if st.button("📸 Summarize Current Screen", key="screen_summary_btn"):
        if not VIDEO_RECORDING_AVAILABLE:
            st.error("❌ Screen capture not available. Install: pip install opencv-python mss")
        elif not subject or subject == "cs101":
            st.error("❌ Please enter Subject and Topic first")
        elif not topic or topic == "lecture demo":
            st.error("❌ Please enter a valid Topic")
        else:
            try:
                frame = capture_screen_frame()
                if frame is None:
                    st.error("❌ Could not capture screen. No monitor detected.")
                else:
                    with st.spinner("📸 Analyzing screen…"):
                        screen_text = ""
                        
                        # Try OCR if available
                        if OCR_AVAILABLE:
                            try:
                                screen_text = pytesseract.image_to_string(frame)
                                if screen_text.strip():
                                    st.success("✅ Text extracted from screen")
                            except Exception as e:
                                st.warning(f"⚠️ OCR error (proceeding without text): {e}")
                        else:
                            st.info("ℹ️ OCR not available - analyzing image only")
                        
                        # Generate prompt for screen analysis
                        prompt = f"""Analyze this screen capture and return ONLY valid JSON.

Subject: {subject} | Topic: {topic}

Screen extracted text: {screen_text[:1000] if screen_text else "(Image analysis only)"}

Return ONLY this JSON structure:
{{
  "screen_summary": "Brief summary of what's displayed on screen",
  "key_content": ["point 1", "point 2", "point 3"],
  "suggested_notes": ["thing to note 1", "thing to note 2"],
  "relevance": 85,
  "recommendations": "Suggestions for lecture notes"
}}
"""
                        
                        try:
                            result = call_ai(prompt, subject, topic)
                            
                            # Try to parse JSON
                            try:
                                parsed = parse_structured_summary(result)
                                st.session_state.summary = json.dumps(parsed)
                                st.success("✅ Screen summary generated!")
                                st.rerun()
                            except Exception as parse_error:
                                st.error(f"❌ Could not parse AI response: {parse_error}")
                                st.write("Raw response:", result[:200])
                        except Exception as ai_error:
                            st.error(f"❌ AI error: {ai_error}")
            except Exception as e:
                st.error(f"❌ Screen capture error: {e}")

    # ----- File Uploads -----
    st.divider()
    st.subheader("📂 Upload Materials")
    uploaded_files = st.file_uploader("Upload PDF or PPTX", type=["pdf", "pptx", "ppt"],
                                      accept_multiple_files=True)
    if uploaded_files:
        for up_file in uploaded_files:
            if up_file.size > MAX_FILE_SIZE:
                st.error(f"❌ File {up_file.name} exceeds 50MB limit")
                continue
            
            file_id = up_file.name + str(up_file.size)
            if file_id not in st.session_state.uploaded_materials:
                st.session_state.uploaded_materials.add(file_id)
                extracted = ""
                if up_file.type == "application/pdf" and PyPDF2:
                    try:
                        reader = PyPDF2.PdfReader(up_file)
                        for page in reader.pages:
                            extracted += page.extract_text() + "\n"
                    except Exception as e:
                        st.error(f"PDF error: {e}")
                elif "presentation" in up_file.type and Presentation:
                    try:
                        with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as tmp:
                            tmp.write(up_file.getvalue())
                            tmp_path = tmp.name
                        try:
                            pres = Presentation(tmp_path)
                            for slide in pres.slides:
                                for shape in slide.shapes:
                                    if hasattr(shape, "text"):
                                        extracted += shape.text + "\n"
                        finally:
                            os.remove(tmp_path)
                    except Exception as e:
                        st.error(f"PPTX error: {e}")
                
                if extracted:
                    st.session_state.material_text += extracted + "\n"
                    st.success(f"✅ Processed {up_file.name}")
        
        if st.button("🗑️ Clear materials"):
            st.session_state.uploaded_materials.clear()
            st.session_state.material_text = ""
            st.rerun()

    # ----- Export / Reset -----
    st.divider()
    st.subheader("📥 Export & Reset")
    if st.session_state.summary:
        st.download_button(
            "💾 Download Summary", 
            st.session_state.summary, 
            f"Summary_{subject}_{session_id}.json", 
            "application/json"
        )
        st.download_button(
            "📜 Download Transcript", 
            st.session_state.transcript, 
            f"Transcript_{subject}_{session_id}.txt", 
            "text/plain"
        )
    
    if st.button("🔄 Reset Session"):
        st.session_state.running = False
        st.session_state.transcript = ""
        st.session_state.summary = ""
        st.session_state.last_summary_len = 0
        st.session_state.student_questions = []
        st.rerun()

# =========================================================
# WHISPER MODEL (cached)
# =========================================================
@st.cache_resource
def load_whisper():
    with st.spinner("⏳ Loading Whisper model (first run only)..."):
        return whisper.load_model("base")

model = load_whisper()

def clear_whisper_memory():
    global model
    try:
        del model
        if torch.cuda.is_available():
            torch.cuda.empty_cache()
    except Exception as e:
        print(f"⚠️ Memory clear error: {e}")

# =========================================================
# AI FUNCTIONS
# =========================================================
def build_prompt(transcript_text, subject, topic):
    """Create a strict JSON prompt for AI summary."""
    context_parts = []
    if st.session_state.material_text:
        context_parts.append(f"--- Lecture Materials ---\n{st.session_state.material_text[:2000]}")
    context_parts.append(f"--- Transcript ---\n{transcript_text}")
    combined = "\n\n".join(context_parts)

    prompt = f"""Analyze this lecture and return ONLY valid JSON.

Use these exact keys and data types:
- "fast_summary": array of 3 short bullet points (strings)
- "normal_explanation": one paragraph string
- "slow_learner_support": simplified explanation string (if needed, otherwise empty string)
- "missed_content": string describing what late students missed (or empty string)
- "clean_summary": overall summary string
- "key_points": array of strings (each a key point)
- "important_concepts": array of strings (each a concept definition)
- "exam_questions": array of at least 3 possible exam questions (strings)
- "revision_notes": array of short revision strings

DO NOT return any other text, only the JSON object.
Subject: {subject} | Topic: {topic}

LECTURE:
{combined}
"""
    return prompt

def call_ai_ollama(prompt):
    try:
        response = ollama.chat(
            model=st.session_state.ollama_model,
            messages=[{"role": "user", "content": prompt}]
        )
        return response['message']['content']
    except Exception as e:
        return json.dumps({"clean_summary": f"❌ Ollama error: {e}"})

def call_ai_groq(prompt, model_to_use=None):
    if not st.session_state.groq_client:
        return None, "Groq client not configured"
    
    if model_to_use is None:
        model_to_use = GROQ_MODEL
    
    try:
        response = st.session_state.groq_client.chat.completions.create(
            model=model_to_use,
            messages=[{"role": "user", "content": prompt}],
            max_tokens=1500
        )
        return response.choices[0].message.content, None
    except Exception as e:
        return None, str(e)

def call_ai_huggingface(prompt, model_to_use=None):
    if not st.session_state.hf_client:
        return None, "Hugging Face client not configured"
    
    if model_to_use is None:
        model_to_use = HF_MODEL
    
    try:
        response = st.session_state.hf_client.chat.completions.create(
            model=model_to_use,
            messages=[{"role": "user", "content": prompt}],
            max_tokens=1500
        )
        return response.choices[0].message.content, None
    except Exception as e:
        return None, str(e)

def call_ai(prompt, subject, topic):
    """Call AI based on selected provider"""
    # For screen summarization, don't rebuild prompt
    if "Analyze this lecture" not in prompt:
        # This is a screen summary prompt, use as-is
        pass
    else:
        # This is a transcript prompt, rebuild it
        prompt = build_prompt(prompt, subject, topic)
    
    # 1. Ollama if selected
    if "Ollama" in st.session_state.ai_provider:
        return call_ai_ollama(prompt)
    
    # 2. Groq if selected
    elif "Groq" in st.session_state.ai_provider:
        result, error = call_ai_groq(prompt)
        if result is not None:
            return result
        st.warning("⚠️ Primary Groq model failed – trying backup models...")
        for backup in GROQ_BACKUPS:
            result, error = call_ai_groq(prompt, backup)
            if result is not None:
                return result
        st.warning("🌐 All Groq models unavailable – falling back to local Ollama.")
        return call_ai_ollama(prompt)
    
    # 3. Hugging Face if selected
    elif "Hugging Face" in st.session_state.ai_provider:
        result, error = call_ai_huggingface(prompt)
        if result is not None:
            return result
        st.warning("⚠️ Primary HF model failed – trying backup models...")
        for backup in HF_BACKUPS:
            result, error = call_ai_huggingface(prompt, backup)
            if result is not None:
                return result
        st.warning("🌐 All Hugging Face models unavailable – falling back to local Ollama.")
        return call_ai_ollama(prompt)
    
    return call_ai_ollama(prompt)

# ---------- JSON Parser ----------
def parse_structured_summary(raw_text):
    default = {
        "fast_summary": ["Could not parse summary."],
        "normal_explanation": "Parsing error.",
        "slow_learner_support": "No support data.",
        "missed_content": "No missed data.",
        "clean_summary": raw_text,
        "key_points": [],
        "important_concepts": [],
        "exam_questions": [],
        "revision_notes": []
    }
    
    try:
        data = json.loads(raw_text)
        for field in ["fast_summary", "key_points", "important_concepts", "exam_questions", "revision_notes"]:
            if field in data:
                if isinstance(data[field], str):
                    data[field] = [data[field]]
                elif not isinstance(data[field], list):
                    data[field] = []
        return data
    except json.JSONDecodeError:
        pass
    
    try:
        match = re.search(r'```(?:json)?\s*(\{.*?\})\s*```', raw_text, re.DOTALL)
        if match:
            return parse_structured_summary(match.group(1))
    except Exception:
        pass
    
    try:
        start = raw_text.find('{')
        end = raw_text.rfind('}') + 1
        if start != -1 and end > start:
            return parse_structured_summary(raw_text[start:end])
    except Exception:
        pass
    
    return default

# =========================================================
# SCREEN & AUDIO
# =========================================================
def capture_screen_frame():
    if not VIDEO_RECORDING_AVAILABLE:
        return None
    try:
        with mss.mss() as sct:
            if len(sct.monitors) < 1:
                return None
            monitor = sct.monitors[1] if len(sct.monitors) > 1 else sct.monitors[0]
            img = sct.grab(monitor)
            pil_img = Image.frombytes("RGB", img.size, img.rgb)
            pil_img.thumbnail((SCREEN_CAPTURE_WIDTH, SCREEN_CAPTURE_HEIGHT))
            return pil_img
    except Exception as e:
        st.session_state.error_msg = f"Screen capture error: {e}"
        return None

def screen_recording_worker(filepath, stop_event):
    out = None
    try:
        with mss.mss() as sct:
            if len(sct.monitors) < 1:
                st.error("No monitor detected")
                return
            monitor = sct.monitors[1] if len(sct.monitors) > 1 else sct.monitors[0]
            fourcc = cv2.VideoWriter_fourcc(*"XVID")
            out = cv2.VideoWriter(filepath, fourcc, 10.0, (monitor["width"], monitor["height"]))
            while not stop_event.is_set():
                img = sct.grab(monitor)
                frame = np.array(img)
                frame = cv2.cvtColor(frame, cv2.COLOR_BGRA2BGR)
                out.write(frame)
                time.sleep(0.1)
    except Exception as e:
        st.error(f"Screen recording error: {e}")
    finally:
        if out:
            out.release()

def audio_transcription_worker(duration, samplerate, stop_event, text_queue, silence_limit):
    while not stop_event.is_set():
        try:
            audio = sd.rec(int(duration * samplerate), samplerate=samplerate, channels=1)
            sd.wait()
        except sd.PortAudioError as e:
            text_queue.put(f" [Audio Device Error: {e}] ")
            time.sleep(1)
            continue
        except Exception as e:
            text_queue.put(f" [Audio Error: {e}] ")
            time.sleep(1)
            continue
        if stop_event.is_set():
            break
        audio_int16 = np.int16(audio * 32767)
        with tempfile.NamedTemporaryFile(delete=False, suffix=".wav") as tmp:
            wav.write(tmp.name, samplerate, audio_int16)
            tmp_path = tmp.name
        try:
            result = model.transcribe(tmp_path, fp16=False)
            speech = result.get("text", "").strip()
            if len(speech) >= silence_limit:
                text_queue.put(speech + " ")
        except Exception as e:
            text_queue.put(f" [Transcription Error: {e}] ")
        finally:
            try:
                os.remove(tmp_path)
            except OSError:
                pass

# =========================================================
# BOTTOM HELPER TABS - MANUAL MODE
# =========================================================
st.divider()
col_notes, col_qa = st.columns(2)

with col_notes:
    st.subheader("💡 AI Helper & Notes")
    
    # ===== MANUAL AI HELPER =====
    st.markdown("**Ask AI (Manual)**")
    notion_prompt = st.text_area("Ask AI (clarification, etc.)", key="notion_prompt", height=80)
    if st.button("🚀 Send to AI", key="send_to_ai_btn"):
        if notion_prompt.strip():
            with st.spinner("🤔 Thinking…"):
                short_prompt = f"Answer concisely in 1-2 sentences:\n{notion_prompt}"
                try:
                    if "Ollama" in st.session_state.ai_provider:
                        response = ollama.chat(
                            model=st.session_state.ollama_model,
                            messages=[{"role": "user", "content": short_prompt}]
                        )
                        st.info(response['message']['content'])
                    elif "Groq" in st.session_state.ai_provider and st.session_state.groq_client:
                        response = st.session_state.groq_client.chat.completions.create(
                            model=GROQ_MODEL,
                            messages=[{"role": "user", "content": short_prompt}],
                            max_tokens=200
                        )
                        st.info(response.choices[0].message.content)
                    elif "Hugging Face" in st.session_state.ai_provider and st.session_state.hf_client:
                        response = st.session_state.hf_client.chat.completions.create(
                            model=HF_MODEL,
                            messages=[{"role": "user", "content": short_prompt}],
                            max_tokens=200
                        )
                        st.info(response.choices[0].message.content)
                    else:
                        st.error("❌ AI provider not configured")
                except Exception as e:
                    st.error(f"❌ AI error: {e}")
    
    # ===== MANUAL SAVE NOTES =====
    st.divider()
    st.markdown("**💾 Saved Notes (Manual Save)**")
    note_title = st.text_input("Note title", key="note_title_input", placeholder="Enter title...")
    note_content = st.text_area("Note content", key="note_content_input", height=80, placeholder="Enter content...")
    if st.button("💾 Save note", key="save_note_btn"):
        if note_title.strip() and note_content.strip():
            if add_note(note_title, note_content):
                st.success("✅ Note saved")
            else:
                st.error("❌ Failed to save note")
        else:
            st.error("❌ Please enter both title and content")
    
    # ===== SEARCH NOTES =====
    st.divider()
    st.markdown("**🔍 Search Notes**")
    search_q = st.text_input("Search notes", key="search_notes_input", placeholder="Search by keyword...")
    if st.button("🔍 Search", key="search_btn"):
        if search_q.strip():
            results = search_notes(search_q)
            if results:
                st.markdown("**Results:**")
                for t, c, cr in results:
                    with st.expander(f"📌 {t} _({cr[:10]})_"):
                        st.write(c)
            else:
                st.info("No notes found")

with col_qa:
    st.subheader("👋 Student Questions")
    
    # ===== MANUAL RAISE HAND =====
    st.markdown("**Submit Question (Manual)**")
    new_question = st.text_input("Student Name", placeholder="Anonymous", key="student_name_input")
    q_text = st.text_area("Question", key="question_input", height=100, placeholder="Enter question...")
    if st.button("✋ Raise Hand", key="raise_hand_btn"):
        if q_text.strip():
            st.session_state.student_questions.append({
                "name": new_question.strip() if new_question.strip() else "Anonymous",
                "text": q_text.strip(),
                "timestamp": datetime.now().isoformat()
            })
            st.success("✅ Question submitted")
            st.rerun()
        else:
            st.error("❌ Please enter a question")
    
    # ===== QUESTIONS QUEUE =====
    if st.session_state.student_questions:
        st.divider()
        st.markdown("**📋 Questions Queue:**")
        for i, q in enumerate(st.session_state.student_questions, 1):
            col_q, col_del = st.columns([5, 1])
            with col_q:
                st.markdown(f"**{i}. {q['name']}:** {q['text']}")
                st.caption(f"⏰ {q['timestamp'][:16]}")
            with col_del:
                if st.button("✅", key=f"del_q_{i}"):
                    st.session_state.student_questions.pop(i - 1)
                    st.rerun()

# =========================================================
# CONTROL BUTTONS
# =========================================================
c1, c2, c3 = st.sidebar.columns(3)
with c1:
    if st.button("▶ Start", use_container_width=True):
        if not subject or subject == "cs101":
            st.error("❌ Enter a valid subject")
        else:
            st.session_state.running = True
            st.session_state.last_summary_len = 0
            st.session_state.recording_stop_event.clear()
            if st.session_state.record_screen and VIDEO_RECORDING_AVAILABLE:
                filepath = os.path.join(REP_DIR, f"{subject}_{session_id}_screen.avi")
                t = threading.Thread(target=screen_recording_worker,
                                     args=(filepath, st.session_state.recording_stop_event))
                t.daemon = True
                t.start()
            if st.session_state.live_audio:
                t2 = threading.Thread(
                    target=audio_transcription_worker,
                    args=(st.session_state.recording_duration, AUDIO_SAMPLE_RATE,
                          st.session_state.recording_stop_event,
                          st.session_state.audio_queue, st.session_state.silence_filter)
                )
                t2.daemon = True
                t2.start()
            st.rerun()
with c2:
    if st.button("⏹ Stop", use_container_width=True):
        st.session_state.running = False
        st.session_state.recording_stop_event.set()
        clear_whisper_memory()
        st.rerun()
with c3:
    if st.button("🔄 Reset", use_container_width=True):
        st.session_state.running = False
        st.session_state.transcript = ""
        st.session_state.summary = ""
        st.session_state.last_summary_len = 0
        st.session_state.student_questions = []
        st.session_state.recording_stop_event.set()
        st.rerun()

if st.session_state.running:
    st.sidebar.success(f"🔴 Recording... (Using {st.session_state.ai_provider})")
else:
    st.sidebar.info("⏸️ Paused")

# =========================================================
# DISPLAY TRANSCRIPT & SUMMARY
# =========================================================
if st.session_state.show_screen and VIDEO_RECORDING_AVAILABLE:
    frame = capture_screen_frame()
    if frame is not None:
        screen_placeholder.image(frame, use_container_width=True)

transcript_display = st.session_state.transcript[-2000:] if st.session_state.transcript else "Waiting to start..."
transcript_placeholder.text_area("Live Transcript", transcript_display, height=350, disabled=True)

if st.session_state.summary:
    structured = parse_structured_summary(st.session_state.summary)
    tabs = summary_placeholder.tabs(
        ["⚡Fast", "📖Normal", "🐢Slow", "⚠️Missed", "📌Summary", "🔑Key", "🧠Concepts", "❓Exam", "📚Revise"]
    )
    with tabs[0]:
        for item in structured.get("fast_summary", []):
            st.write(f"• {item}")
    with tabs[1]:
        st.write(structured.get("normal_explanation", ""))
    with tabs[2]:
        st.write(structured.get("slow_learner_support", ""))
    with tabs[3]:
        st.write(structured.get("missed_content", ""))
    with tabs[4]:
        st.write(structured.get("clean_summary", ""))
    with tabs[5]:
        for item in structured.get("key_points", []):
            st.write(f"• {item}")
    with tabs[6]:
        for item in structured.get("important_concepts", []):
            st.write(f"• {item}")
    with tabs[7]:
        for i, q in enumerate(structured.get("exam_questions", []), 1):
            st.write(f"{i}. {q}")
    with tabs[8]:
        for item in structured.get("revision_notes", []):
            st.write(f"• {item}")
else:
    summary_placeholder.info("📊 Summary will appear here.")

# =========================================================
# MANUAL SUMMARY BUTTON
# =========================================================
st.divider()
st.subheader("📊 Manual Summary Generation")
col_gen_left, col_gen_right = st.columns([3, 1])
with col_gen_left:
    if st.session_state.transcript:
        st.write(f"**Transcript length:** {len(st.session_state.transcript)} chars")
with col_gen_right:
    if st.button("🚀 Generate Summary Now", use_container_width=True):
        if not st.session_state.transcript.strip():
            st.error("❌ No transcript available.")
        else:
            with st.spinner("🧠 Generating AI summary..."):
                new_summary = call_ai(st.session_state.transcript, subject, topic)
                st.session_state.summary = new_summary
                st.session_state.last_summary_len = len(st.session_state.transcript)
                filepath = os.path.join(REP_DIR, f"{subject}_{session_id}.json")
                try:
                    with open(filepath, "w", encoding="utf-8") as f:
                        f.write(new_summary)
                    st.success(f"✅ Summary saved!")
                except Exception as e:
                    st.error(f"❌ Error saving: {e}")
                st.rerun()

# =========================================================
# MAIN PROCESSING LOOP
# =========================================================
if st.session_state.running:
    new_text_added = False
    while not st.session_state.audio_queue.empty():
        text = st.session_state.audio_queue.get()
        st.session_state.transcript += text
        new_text_added = True
    if new_text_added:
        st.session_state.transcript = st.session_state.transcript[-8000:]

    transcript_len = len(st.session_state.transcript)
    if (transcript_len > st.session_state.summary_trigger and 
        transcript_len - st.session_state.last_summary_len > 200):
        with st.spinner("🧠 Generating AI summary..."):
            new_summary = call_ai(st.session_state.transcript, subject, topic)
            st.session_state.summary = new_summary
            st.session_state.last_summary_len = transcript_len
            filepath = os.path.join(REP_DIR, f"{subject}_{session_id}.json")
            try:
                with open(filepath, "w", encoding="utf-8") as f:
                    f.write(new_summary)
            except Exception as e:
                st.session_state.error_msg = f"Save error: {e}"
    time.sleep(0.5)
    st.rerun()
