import streamlit as st
st.set_page_config(page_title="AI Lecture Assistant", page_icon="🧠", layout="wide")

import os
try:
    import whisper
    WHISPER_AVAILABLE = True
except ImportError:
    WHISPER_AVAILABLE = False

try:
    import sounddevice as sd
    import scipy.io.wavfile as wav
    AUDIO_AVAILABLE = True
except ImportError:
    AUDIO_AVAILABLE = False

import tempfile
import json
import time
import re
import contextlib
import threading
import queue
import numpy as np
import torch
from datetime import datetime
from sqlite3 import connect as sqlite_connect
from PIL import Image
import base64
import io
import pandas as pd

try:
    import cv2
    import mss
    VIDEO_RECORDING_AVAILABLE = True
except ImportError:
    VIDEO_RECORDING_AVAILABLE = False

try:
    import PyPDF2
except ImportError:
    PyPDF2 = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

try:
    import ollama
    OLLAMA_AVAILABLE = True
except ImportError:
    OLLAMA_AVAILABLE = False

try:
    from groq import Groq as GroqClient
    GROQ_AVAILABLE = True
except ImportError:
    GROQ_AVAILABLE = False

try:
    from huggingface_hub import InferenceClient as HFInferenceClient
    HF_AVAILABLE = True
except ImportError:
    HF_AVAILABLE = False

# ---------- Optional OCR ----------
try:
    import pytesseract
    import shutil
    tesseract_path = shutil.which("tesseract")
    if tesseract_path:
        pytesseract.pytesseract.tesseract_cmd = tesseract_path
    elif os.path.exists(r"C:\Program Files\Tesseract-OCR\tesseract.exe"):
        pytesseract.pytesseract.tesseract_cmd = r"C:\Program Files\Tesseract-OCR\tesseract.exe"
    OCR_AVAILABLE = True
except ImportError:
    OCR_AVAILABLE = False

# =========================================================
# CONSTANTS AND CONFIG
# =========================================================
BASE_DIR = "app_data"
REP_DIR = os.path.join(BASE_DIR, "reports")
ATTENDANCE_FILE = os.path.join(BASE_DIR, "attendance.csv")
os.makedirs(REP_DIR, exist_ok=True)

OLLAMA_MODEL = "llama3"
GROQ_MODEL = "llama-3.3-70b-versatile"
GROQ_BACKUPS = ["qwen-2.5-32b", "deepseek-r1-distill-llama-70b"]
HF_MODEL = "mistralai/Mistral-7B-Instruct-v0.3"
HF_BACKUPS = ["Qwen/Qwen2.5-7B-Instruct", "meta-llama/Llama-3.3-70B-Instruct"]

AUDIO_SAMPLE_RATE = 16000
SCREEN_CAPTURE_WIDTH = 800
SCREEN_CAPTURE_HEIGHT = 600

# =========================================================
# SESSION STATE INITIALIZATION
# =========================================================
def init_session_state():
    defaults = {
        "running": False,
        "transcript": "",
        "summary": "",
        "last_summary_len": 0,
        "error_msg": "",
        "uploaded_materials": set(),
        "material_text": "",
        "last_vision_time": 0,
        "recording_stop_event": threading.Event(),
        "ai_provider": "Ollama",
        "attendance_data": [],
        "audio_queue": queue.Queue(),
        "history_loaded": False,
        "summary_queue": queue.Queue(),
        "generating_summary": False,
        "cancel_summary_event": threading.Event(),
    }
    for key, default_value in defaults.items():
        if key not in st.session_state:
            st.session_state[key] = default_value

init_session_state()

st.title("🧠 AI Lecture Assistant — Professor Edition")

if st.session_state.error_msg:
    st.warning(st.session_state.error_msg)
    st.session_state.error_msg = ""

# =========================================================
#  ALL HELPER FUNCTIONS DEFINED HERE
# =========================================================

# ---- Database ----
@contextlib.contextmanager
def get_db_connection():
    db_path = os.path.join(BASE_DIR, "notes.db")
    conn = sqlite_connect(db_path)
    try:
        yield conn
    finally:
        conn.close()

def init_notes_db():
    with get_db_connection() as conn:
        conn.execute(
            "CREATE TABLE IF NOT EXISTS notes ("
            "id INTEGER PRIMARY KEY AUTOINCREMENT, "
            "title TEXT NOT NULL, "
            "content TEXT NOT NULL, "
            "created_at TEXT NOT NULL)"
        )
        conn.commit()

init_notes_db()

# ---- Attendance ----
def init_attendance_file():
    if not os.path.exists(ATTENDANCE_FILE):
        df = pd.DataFrame(columns=["Date", "Time", "Name", "Status", "AI_Support"])
        df.to_csv(ATTENDANCE_FILE, index=False)

def mark_attendance(name, status="On-Time"):
    now = datetime.now()
    support = "AI Summary Required" if status == "Late" else "Optional"
    data = {
        "Date": [now.strftime("%Y-%m-%d")],
        "Time": [now.strftime("%H:%M:%S")],
        "Name": [name],
        "Status": [status],
        "AI_Support": [support]
    }
    df = pd.DataFrame(data)
    df.to_csv(ATTENDANCE_FILE, mode='a', index=False, header=False)

def get_latecomers():
    init_attendance_file()
    try:
        df = pd.read_csv(ATTENDANCE_FILE)
        today = datetime.now().strftime("%Y-%m-%d")
        late_mask = (df['Date'] == today) & (df['Status'] == 'Late')
        return df[late_mask]['Name'].unique().tolist()
    except:
        return []

init_attendance_file()

# ---- Whisper ----
@st.cache_resource
def load_whisper():
    if not WHISPER_AVAILABLE:
        return None
    with st.spinner("⏳ Loading Whisper model..."):
        try:
            return whisper.load_model("base")
        except Exception as e:
            st.error(f"Error loading Whisper: {e}")
            return None

model = load_whisper()

def clear_whisper_memory():
    global model
    try:
        del model
        if torch.cuda.is_available():
            torch.cuda.empty_cache()
    except:
        pass

# ---- Screen Capture ----
def capture_screen_frame():
    if not VIDEO_RECORDING_AVAILABLE:
        return None
    try:
        with mss.mss() as sct:
            monitor = sct.monitors[1] if len(sct.monitors) > 1 else sct.monitors[0]
            img = sct.grab(monitor)
            pil_img = Image.frombytes("RGB", img.size, img.rgb)
            pil_img.thumbnail((SCREEN_CAPTURE_WIDTH, SCREEN_CAPTURE_HEIGHT))
            return pil_img
    except Exception as e:
        st.session_state.error_msg = f"Screen capture error: {e}"
        return None

def screen_recording_worker(filepath, stop_event):
    out = None
    try:
        with mss.mss() as sct:
            monitor = sct.monitors[1] if len(sct.monitors) > 1 else sct.monitors[0]
            fourcc = cv2.VideoWriter_fourcc(*"XVID")
            out = cv2.VideoWriter(filepath, fourcc, 10.0, (monitor["width"], monitor["height"]))
            while not stop_event.is_set():
                img = sct.grab(monitor)
                frame = np.array(img)
                frame = cv2.cvtColor(frame, cv2.COLOR_BGRA2BGR)
                out.write(frame)
                time.sleep(0.1)
    except:
        pass
    finally:
        if out:
            out.release()

# ---- Audio ----
def audio_transcription_worker(duration, samplerate, stop_event, text_queue, silence_limit):
    while not stop_event.is_set():
        try:
            audio = sd.rec(int(duration * samplerate), samplerate=samplerate, channels=1)
            sd.wait()
        except:
            time.sleep(1)
            continue
        if stop_event.is_set():
            break
        audio_int16 = np.int16(audio * 32767)
        with tempfile.NamedTemporaryFile(delete=False, suffix=".wav") as tmp:
            wav.write(tmp.name, samplerate, audio_int16)
            tmp_path = tmp.name
        try:
            if model:
                result = model.transcribe(tmp_path, fp16=False)
                speech = result.get("text", "").strip()
                if len(speech) >= silence_limit:
                    text_queue.put(speech + " ")
            else:
                time.sleep(1)
        except:
            pass
        finally:
            try:
                os.remove(tmp_path)
            except OSError:
                pass

# ---- AI Functions ----
def build_prompt(transcript_text):
    context_parts = []
    if st.session_state.material_text:
        context_parts.append(f"--- Lecture Materials ---\n{st.session_state.material_text[:2000]}")
    context_parts.append(f"--- Transcript ---\n{transcript_text}")
    combined = "\n\n".join(context_parts)

    prompt = f"""Analyze this lecture and return ONLY valid JSON.
Use these exact keys and data types:
- "fast_summary": array of 3 short bullet points (strings)
- "normal_explanation": one paragraph string
- "slow_learner_support": simplified explanation string (or empty)
- "missed_content": string for late students (or empty)
- "clean_summary": overall summary string
- "key_points": array of strings
- "important_concepts": array of strings
- "exam_questions": array of at least 3 strings
- "revision_notes": array of strings

DO NOT return any other text, only the JSON object.
Subject: {subject} | Topic: {topic}

LECTURE:
{combined}
"""
    return prompt

def call_ai_ollama(prompt):
    if not OLLAMA_AVAILABLE:
        return json.dumps({"clean_summary": "❌ Ollama library not installed."})
    try:
        model_name = st.session_state.get("ollama_model", OLLAMA_MODEL)
        response = ollama.chat(model=model_name, messages=[{"role": "user", "content": prompt}])
        content = response['message']['content']
        return content
    except Exception as e:
        return json.dumps({"clean_summary": f"❌ Ollama error: {e}. Is Ollama running?"})

def call_ai_groq(prompt, model_to_use=None):
    if not GROQ_AVAILABLE or not groq_client:
        return None, "Groq client not configured or installed"
    model_to_use = model_to_use or GROQ_MODEL
    try:
        response = groq_client.chat.completions.create(
            model=model_to_use,
            messages=[{"role": "user", "content": prompt}],
            max_tokens=1500
        )
        content = response.choices[0].message.content
        return content, None
    except Exception as e:
        return None, str(e)

def call_ai_huggingface(prompt, model_to_use=None):
    if not HF_AVAILABLE or not hf_client:
        return None, "Hugging Face client not configured or installed"
    model_to_use = model_to_use or HF_MODEL
    try:
        response = hf_client.chat.completions.create(
            model=model_to_use,
            messages=[{"role": "user", "content": prompt}],
            max_tokens=1500
        )
        content = response.choices[0].message.content
        return content, None
    except Exception as e:
        return None, str(e)

def call_ai(prompt):
    if "Ollama" in st.session_state.ai_provider:
        return call_ai_ollama(prompt)
    elif "Groq" in st.session_state.ai_provider:
        result, _ = call_ai_groq(prompt)
        if result: return result
        for backup in GROQ_BACKUPS:
            result, _ = call_ai_groq(prompt, backup)
            if result: return result
        return call_ai_ollama(prompt)
    elif "Hugging Face" in st.session_state.ai_provider:
        result, _ = call_ai_huggingface(prompt)
        if result: return result
        for backup in HF_BACKUPS:
            result, _ = call_ai_huggingface(prompt, backup)
            if result: return result
        return call_ai_ollama(prompt)

# ---- background summary worker ----
def summary_background_worker(transcript):
    if st.session_state.cancel_summary_event.is_set():
        st.session_state.summary_queue.put(None)
        return
    prompt = build_prompt(transcript)
    result = call_ai(prompt)
    if not st.session_state.cancel_summary_event.is_set():
        st.session_state.summary_queue.put(result)

def generate_summary_in_background(transcript):
    if st.session_state.generating_summary:
        return
    st.session_state.cancel_summary_event.clear()
    st.session_state.generating_summary = True
    thread = threading.Thread(target=summary_background_worker,
                              args=(transcript,), daemon=True)
    thread.start()

def parse_structured_summary(raw_text):
    default = {
        "fast_summary": ["Could not parse summary."],
        "normal_explanation": "Parsing error.",
        "slow_learner_support": "No support data.",
        "missed_content": "No missed data.",
        "clean_summary": raw_text,
        "key_points": [],
        "important_concepts": [],
        "exam_questions": [],
        "revision_notes": []
    }
    clean = raw_text
    try:
        data = json.loads(clean)
        for field in ["fast_summary", "key_points", "important_concepts", "exam_questions", "revision_notes"]:
            if field in data:
                if isinstance(data[field], str):
                    data[field] = [data[field]]
                elif not isinstance(data[field], list):
                    data[field] = []
        return data
    except json.JSONDecodeError:
        pass
    try:
        match = re.search(r'```(?:json)?\s*(\{.*?\})\s*```', clean, re.DOTALL)
        if match: return parse_structured_summary(match.group(1))
    except: pass
    try:
        start = clean.find('{')
        end = clean.rfind('}') + 1
        if start != -1 and end > start: return parse_structured_summary(clean[start:end])
    except: pass
    return default

# =========================================================
#  SIDEBAR
# =========================================================
groq_client = None
hf_client = None

with st.sidebar:
    st.divider()
    st.subheader("🖥️ System Status")
    status_cols = st.columns(2)
    with status_cols[0]:
        st.write("🎤 Audio:", "✅" if AUDIO_AVAILABLE else "❌")
        st.write("🧠 Whisper:", "✅" if WHISPER_AVAILABLE else "❌")
        st.write("🤖 Ollama:", "✅" if OLLAMA_AVAILABLE else "❌")
    with status_cols[1]:
        st.write("📸 Screen:", "✅" if VIDEO_RECORDING_AVAILABLE else "❌")
        st.write("🔍 OCR:", "✅" if OCR_AVAILABLE else "❌")
        st.write("☁️ Groq:", "✅" if GROQ_AVAILABLE else "❌")
    
    if not WHISPER_AVAILABLE or not AUDIO_AVAILABLE:
        st.error("Missing critical libraries. Run: pip install openai-whisper sounddevice scipy")

    st.divider()
    st.subheader("🤖 AI Provider Setup")
    ai_provider = st.radio(
        "Select AI Provider",
        ["Ollama (Local, Free)", "Groq (Cloud, Fast)", "Hugging Face (Cloud)"],
        index=0
    )
    st.session_state.ai_provider = ai_provider.split(" ")[0]

    if "Ollama" in ai_provider:
        if not OLLAMA_AVAILABLE:
            st.error("❌ Ollama not installed.")
        else:
            available_models = [
                "llama3:latest",
                "phi4-mini:latest",
                "llama3.1:latest",
                "phi3:latest",
                "qwen2.5:3b",
                "deepseek-r1:1.5b",
                "Custom..."
            ]
            selected_model = st.selectbox("Ollama Model", available_models, index=0)
            if selected_model == "Custom...":
                custom_model = st.text_input("Enter custom model name", value=OLLAMA_MODEL)
                st.session_state.ollama_model = custom_model
            else:
                st.session_state.ollama_model = selected_model
    
    elif "Groq" in ai_provider:
        if not GROQ_AVAILABLE:
            st.error("❌ Groq not installed.")
        else:
            groq_key = st.text_input("🔑 Groq API Key", type="password", placeholder="gsk_...")
            if groq_key:
                try:
                    groq_client = GroqClient(api_key=groq_key)
                    st.success("✅ Groq connected")
                except Exception as e:
                    st.error(f"❌ Groq error: {e}")
    elif "Hugging Face" in ai_provider:
        if not HF_AVAILABLE:
            st.error("❌ Hugging Face hub not installed.")
        else:
            hf_token = st.text_input("🤗 HF Token", type="password", placeholder="hf_...")
            if hf_token:
                try:
                    hf_client = HFInferenceClient(token=hf_token)
                    st.success("✅ Hugging Face connected")
                except Exception as e:
                    st.error(f"❌ HF error: {e}")

    st.divider()
    st.header("⚙️ Session Controls")
    subject = st.text_input("📘 Subject", "CS101").strip().lower().replace(" ", "_")
    session_id = st.text_input("🆔 Session ID", datetime.now().strftime("%Y%m%d_%H%M"))
    topic = st.text_input("📚 Topic", "Lecture Demo")

    st.divider()
    live_audio = st.checkbox("🎤 Live Voice Capture", value=True)
    show_screen = st.checkbox("📺 Show Presentation Screen", value=True, disabled=not VIDEO_RECORDING_AVAILABLE)
    record_screen = st.checkbox("🎥 Record Screen to Video", value=False, disabled=not VIDEO_RECORDING_AVAILABLE)
    
    silence_filter = st.slider("🔇 Minimum chars to transcribe", 0, 50, 5)
    summary_trigger = st.slider("Summary after (chars)", 100, 1000, 400)
    recording_duration = st.slider("🎙️ Recording duration (sec)", 2, 10, 5)

    # ----- One-Click Screen Summary -----
    st.divider()
    st.subheader("🖥️ Quick Screen Summary")
    if st.button("Summarize Current Screen"):
        frame = capture_screen_frame()
        if frame:
            with st.spinner("Analysing screen…"):
                screen_text = ""
                if OCR_AVAILABLE:
                    try:
                        screen_text = pytesseract.image_to_string(frame)
                    except:
                        pass
                prompt = f"""Analyze the current screen content.
Return a concise summary in JSON format.
Screen text: {screen_text[:2000]}
Subject: {subject} | Topic: {topic}
"""
                result = call_ai(prompt)
                st.session_state.summary = result
                st.success("Screen summary generated!")
        else:
            st.error("Screen capture not available.")

    # ----- File Uploads -----
    st.divider()
    st.subheader("📂 Upload Materials")
    uploaded_files = st.file_uploader("Upload PDF or PPTX", type=["pdf", "pptx", "ppt"], accept_multiple_files=True)
    if uploaded_files:
        for up_file in uploaded_files:
            file_id = up_file.name + str(up_file.size)
            if file_id not in st.session_state.uploaded_materials:
                st.session_state.uploaded_materials.add(file_id)
                extracted = ""
                if up_file.type == "application/pdf" and PyPDF2:
                    try:
                        reader = PyPDF2.PdfReader(up_file)
                        for page in reader.pages:
                            extracted += page.extract_text() + "\n"
                    except Exception as e:
                        st.error(f"PDF error: {e}")
                elif "presentation" in up_file.type and Presentation:
                    try:
                        with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as tmp:
                            tmp.write(up_file.getvalue())
                            tmp_path = tmp.name
                        try:
                            pres = Presentation(tmp_path)
                            for slide in pres.slides:
                                for shape in slide.shapes:
                                    if hasattr(shape, "text"):
                                        extracted += shape.text + "\n"
                        finally:
                            os.remove(tmp_path)
                    except Exception as e:
                        st.error(f"PPTX error: {e}")
                if extracted:
                    st.session_state.material_text += extracted + "\n"
                    st.success(f"✅ Processed {up_file.name}")
        if st.button("🗑️ Clear materials"):
            st.session_state.uploaded_materials.clear()
            st.session_state.material_text = ""
            st.rerun()

    # ----- Export & Reset -----
    st.divider()
    st.subheader("📥 Export & Reset")
    if st.session_state.summary:
        st.download_button("💾 Download Summary", st.session_state.summary, f"Summary_{subject}_{session_id}.json", "application/json")
    if st.button("🔄 Reset Session"):
        st.session_state.running = False
        st.session_state.transcript = ""
        st.session_state.summary = ""
        st.session_state.last_summary_len = 0
        st.rerun()

# =========================================================
# UI LAYOUT
# =========================================================
col_screen, col_transcript, col_summary = st.columns([1.2, 1, 1])
with col_screen:
    st.subheader("📺 Presentation View")
    screen_placeholder = st.empty()
with col_transcript:
    st.subheader("📄 Live Transcript")
    transcript_placeholder = st.empty()
with col_summary:
    st.subheader("🧠 AI Summary")
    summary_placeholder = st.empty()

# =========================================================
# CONTROL BUTTONS
# =========================================================
c1, c2, c3 = st.sidebar.columns(3)
with c1:
    if st.button("▶ Start"):
        if not subject:
            st.error("❌ Enter subject")
        else:
            st.session_state.running = True
            st.session_state.last_summary_len = 0
            st.session_state.recording_stop_event.clear()
            if record_screen and VIDEO_RECORDING_AVAILABLE:
                filepath = os.path.join(REP_DIR, f"{subject}_{session_id}_screen.avi")
                t = threading.Thread(target=screen_recording_worker, args=(filepath, st.session_state.recording_stop_event))
                t.daemon = True
                t.start()
            if live_audio:
                t2 = threading.Thread(target=audio_transcription_worker,
                                     args=(recording_duration, AUDIO_SAMPLE_RATE,
                                           st.session_state.recording_stop_event,
                                           st.session_state.audio_queue, silence_filter))
                t2.daemon = True
                t2.start()
            st.rerun()
with c2:
    if st.button("⏹ Stop"):
        st.session_state.running = False
        st.session_state.recording_stop_event.set()
        clear_whisper_memory()
        st.rerun()
with c3:
    if st.button("🔄 Reset"):
        st.session_state.running = False
        st.session_state.transcript = ""
        st.session_state.summary = ""
        st.session_state.last_summary_len = 0
        st.session_state.recording_stop_event.set()
        st.rerun()

if st.session_state.running:
    st.sidebar.success(f"🔴 Recording... (Using {st.session_state.ai_provider})")
else:
    st.sidebar.info("⏸️ Paused")

# =========================================================
# MAIN PROCESSING LOOP
# =========================================================
if st.session_state.running:
    new_text_added = False
    while not st.session_state.audio_queue.empty():
        text = st.session_state.audio_queue.get()
        st.session_state.transcript += text
        new_text_added = True
    if new_text_added:
        st.session_state.transcript = st.session_state.transcript[-8000:]

    transcript_len = len(st.session_state.transcript)
    if (transcript_len > summary_trigger and transcript_len - st.session_state.last_summary_len > 200):
        with st.spinner("🧠 Generating AI summary..."):
            new_summary = call_ai(st.session_state.transcript)
            st.session_state.summary = new_summary
            st.session_state.last_summary_len = transcript_len
            report_path = os.path.join(REP_DIR, f"{subject}_{session_id}.json")
            with open(report_path, "w", encoding="utf-8") as f:
                json.dump({"summary": new_summary, "transcript": st.session_state.transcript}, f)

# =========================================================
# DISPLAY TRANSCRIPT & SUMMARY
# =========================================================
if show_screen and VIDEO_RECORDING_AVAILABLE:
    frame = capture_screen_frame()
    if frame:
        screen_placeholder.image(frame, use_container_width=True)

transcript_display = st.session_state.transcript[-2000:] if st.session_state.transcript else "Waiting to start..."
transcript_placeholder.text_area("Live Transcript", transcript_display, height=350, disabled=True)

if st.session_state.summary:
    structured = parse_structured_summary(st.session_state.summary)
    tabs = summary_placeholder.tabs(
        ["⚡Fast", "📖Normal", "🐢Slow", "⚠️Missed", "📌Summary", "🔑Key", "🧠Concepts", "❓Exam", "📚Revise"]
    )
    with tabs[0]:
        for item in structured.get("fast_summary", []): st.write(f"• {item}")
    with tabs[1]: st.write(structured.get("normal_explanation", ""))
    with tabs[2]: st.write(structured.get("slow_learner_support", ""))
    with tabs[3]: st.write(structured.get("missed_content", ""))
    with tabs[4]: st.write(structured.get("clean_summary", ""))
    with tabs[5]:
        for item in structured.get("key_points", []): st.write(f"• {item}")
    with tabs[6]:
        for item in structured.get("important_concepts", []): st.write(f"• {item}")
    with tabs[7]:
        for i, q in enumerate(structured.get("exam_questions", []), 1): st.write(f"{i}. {q}")
    with tabs[8]:
        for item in structured.get("revision_notes", []): st.write(f"• {item}")
else:
    summary_placeholder.info("📊 Summary will appear here.")

# =========================================================
# MANUAL SUMMARY BUTTON
# =========================================================
st.divider()
st.subheader("📊 Manual Summary Generation")
col_gen_left, col_gen_right = st.columns([3, 1])
with col_gen_left:
    if st.session_state.transcript:
        st.write(f"**Transcript length:** {len(st.session_state.transcript)} chars")
with col_gen_right:
    if st.button("🚀 Generate Summary Now"):
        if not st.session_state.transcript.strip():
            st.error("❌ No transcript available.")
        else:
            with st.spinner("🧠 Generating AI summary..."):
                new_summary = call_ai(st.session_state.transcript)
                st.session_state.summary = new_summary
                st.session_state.last_summary_len = len(st.session_state.transcript)
                filepath = os.path.join(REP_DIR, f"{subject}_{session_id}.json")
                with open(filepath, "w", encoding="utf-8") as f:
                    json.dump({"summary": new_summary, "transcript": st.session_state.transcript}, f)
                st.success("✅ Summary saved!")
                st.rerun()

# =========================================================
# FINAL RERUN FOR LIVE UPDATES
# =========================================================
if st.session_state.running:
    time.sleep(1)
    st.rerun()