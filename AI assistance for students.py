import streamlit as st
import whisper
import tempfile
import os
import json
import ollama
import PyPDF2
import pytesseract
import gc
import io
from PIL import Image
from pptx import Presentation
from datetime import datetime


# ---------------- 1. CORE AI ENGINE ----------------
def call_ollama_safe(prompt, model_name, is_json=True):
    """
    is_json=True: Forces JSON output (for study notes).
    is_json=False: Allows natural conversation (for Doubt Solver).
    """
    try:
        response = ollama.generate(
            model=model_name,
            prompt=prompt,
            format="json" if is_json else "",
            options={"num_ctx": 4096, "temperature": 0.3}
        )
        return response['response']
    except Exception as e:
        st.error(f"Connection Error: {e}")
        return None


def parse_json_safely(raw_text):
    try:
        return json.loads(raw_text)
    except:
        return {}


# ---------------- 2. UI SETUP ----------------
st.set_page_config(page_title="Student AI - Simple Edition", layout="wide")
st.title("📚 Student AI Assistant – Study Mode")

if "transcript" not in st.session_state: st.session_state.transcript = ""
if "summary_data" not in st.session_state: st.session_state.summary_data = {}

with st.sidebar:
    st.header("⚙️ Configuration")
    model_choice = st.selectbox("Select AI Brain", ["qwen2.5:3b", "deepseek-r1:1.5b", "phi4-mini:latest"])
    subject = st.text_input("📘 Subject", "Professional Ethics")
    topic = st.text_input("📚 Topic", "General")

    if st.button("🗑️ Clear Session & RAM"):
        st.session_state.transcript = ""
        st.session_state.summary_data = {}
        gc.collect()
        st.rerun()

# ---------------- 3. SIMPLIFIED FILE UPLOAD ----------------
# Removed the "Record" tab as requested. Just a clean file uploader.
uploaded_file = st.file_uploader("Upload Materials (PDF, PPTX, Image)",
                                 type=["pdf", "pptx", "png", "jpg"])

# ---------------- 4. PROCESSING LOGIC ----------------
if st.button("🚀 Generate Detailed Study Notes") and uploaded_file:
    with st.status("Analyzing content...", expanded=True) as status:
        text = ""
        ext = uploaded_file.name.split('.')[-1].lower()
        content = uploaded_file.read()

        # Extracting Text
        if ext == "pdf":
            pdf = PyPDF2.PdfReader(io.BytesIO(content))
            text = " ".join([p.extract_text() for p in pdf.pages])
        elif ext in ["pptx", "ppt"]:
            prs = Presentation(io.BytesIO(content))
            text = " ".join([s.text for slide in prs.slides for s in slide.shapes if hasattr(s, "text")])
        elif ext in ["png", "jpg"]:
            text = pytesseract.image_to_string(Image.open(io.BytesIO(content)))

        if text.strip():
            st.session_state.transcript = text
            status.write(f"🧠 {model_choice} is analyzing...")

            prompt = f"""
            Analyze this {subject} content on {topic}. 
            Return ONLY JSON with these exact keys: 
            "thinking", "explanation", "highlights", "missed", "exam".
            Text: {text[:3000]}
            """

            raw_json = call_ollama_safe(prompt, model_choice, is_json=True)
            st.session_state.summary_data = parse_json_safely(raw_json)
            status.update(label="✅ Analysis Complete!", state="complete")
        else:
            st.error("Could not extract any text.")

# ---------------- 5. DISPLAY & DOWNLOAD ----------------
if st.session_state.summary_data:
    data = st.session_state.summary_data

    # Reasoning Expandable
    if data.get("thinking"):
        with st.expander("💭 View AI's Reasoning Process"):
            st.write(data["thinking"])

    col1, col2 = st.columns([2, 1])
    with col1:
        st.subheader("📖 Full Explanation")
        st.write(data.get("explanation", "N/A"))

    with col2:
        st.subheader("⚡ Key Highlights")
        for h in data.get("highlights", []): st.write(f"- {h}")
        st.success("**Potential Exam Questions**")
        for q in data.get("exam", []): st.write(f"❓ {q}")

    # Feature: Download Study Pack
    st.divider()
    st.download_button(
        label="📥 Download Study Pack (JSON)",
        data=json.dumps(data, indent=4),
        file_name=f"{topic}_notes.json",
        mime="application/json"
    )

    # UPDATED DOUBT SOLVER: No more JSON format errors
    st.subheader("❓ Quick Doubt Solver")
    user_query = st.text_input("Ask a question about this material (e.g., 'What is Utility?'):")

    if user_query:
        with st.spinner("AI is typing..."):
            # is_json=False allows the AI to answer in plain text!
            ans = call_ollama_safe(
                f"Context: {st.session_state.transcript[:2000]}\nQuestion: {user_query}",
                model_choice,
                is_json=False
            )
            st.info(f"**Tutor Answer:** {ans}")
